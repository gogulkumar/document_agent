"""task_ppt_export — generate a richer PowerPoint presentation using python-pptx."""
from __future__ import annotations
import json
import logging
import os
import uuid

from tools.tasks.task_base import invoke_task_llm, EXPORT_ROOT

logger = logging.getLogger(__name__)

_PPT_PROMPT = """
Generate a PowerPoint outline as JSON with this structure:
{
  "title": "Presentation Title",
  "subtitle": "Short positioning statement",
  "theme": {
    "accent": "#1D4ED8",
    "background": "#F8FAFC",
    "surface": "#E2E8F0",
    "text": "#0F172A"
  },
  "slides": [
    {
      "layout": "title|summary|mind_map|brainstorm|table",
      "title": "Slide Title",
      "bullets": ["Point 1", "Point 2"],
      "notes": "Speaker notes for this slide",
      "summary": "Short slide takeaway",
      "nodes": ["Center idea", "Related idea A", "Related idea B"],
      "connections": [["Center idea", "Related idea A"], ["Center idea", "Related idea B"]],
      "ideas": ["Idea 1", "Idea 2", "Idea 3"],
      "table": {
        "headers": ["Column 1", "Column 2"],
        "rows": [["A", "B"], ["C", "D"]]
      }
    }
  ]
}
Respond ONLY with valid JSON. No markdown wrapping.
Make the deck strategic and visual, not just a dump of bullets.
If the content is analytical, include at least one "mind_map" or "brainstorm" slide.
"""


def _hex_to_rgb(value: str, fallback: tuple[int, int, int]) -> tuple[int, int, int]:
    value = (value or "").strip().lstrip("#")
    if len(value) != 6:
        return fallback
    try:
        return tuple(int(value[i:i + 2], 16) for i in (0, 2, 4))
    except Exception:
        return fallback


def _apply_text_style(paragraph, size: int = 20, bold: bool = False, rgb: tuple[int, int, int] = (15, 23, 42)):
    from pptx.util import Pt  # type: ignore
    from pptx.dml.color import RGBColor  # type: ignore

    if paragraph.runs:
        run = paragraph.runs[0]
    else:
        run = paragraph.add_run()
        run.text = paragraph.text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*rgb)


def _add_textbox(slide, left, top, width, height, text, *, size=20, bold=False, color=(15, 23, 42)):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    _apply_text_style(p, size=size, bold=bold, rgb=color)
    return textbox


def _add_bullet_block(slide, left, top, width, height, bullets, *, color=(15, 23, 42)):
    from pptx.util import Pt  # type: ignore
    from pptx.dml.color import RGBColor  # type: ignore

    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, bullet in enumerate(bullets or []):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        if p.runs:
            run = p.runs[0]
        else:
            run = p.add_run()
            run.text = bullet
        run.font.size = Pt(20)
        run.font.color.rgb = RGBColor(*color)
    return textbox


def _add_card(slide, left, top, width, height, title, body, *, fill_rgb, text_rgb):
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE  # type: ignore
    from pptx.dml.color import RGBColor  # type: ignore
    from pptx.util import Inches

    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill_rgb)
    shape.line.color.rgb = RGBColor(*fill_rgb)
    _add_textbox(slide, left + Inches(0.12), top + Inches(0.08), width - Inches(0.24), Inches(0.35), title, size=18, bold=True, color=text_rgb)
    _add_textbox(slide, left + Inches(0.12), top + Inches(0.42), width - Inches(0.24), height - Inches(0.5), body, size=15, color=text_rgb)
    return shape


def _render_mind_map(slide, slide_data, *, accent_rgb, surface_rgb, text_rgb):
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR  # type: ignore
    from pptx.dml.color import RGBColor  # type: ignore
    from pptx.util import Inches

    nodes = slide_data.get("nodes") or []
    center_label = nodes[0] if nodes else slide_data.get("summary") or slide_data.get("title", "Core Idea")
    outer_nodes = nodes[1:] if len(nodes) > 1 else slide_data.get("bullets", [])[:4]
    center = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(3.6), Inches(2.0), Inches(2.0), Inches(2.0))
    center.fill.solid()
    center.fill.fore_color.rgb = RGBColor(*accent_rgb)
    center.line.color.rgb = RGBColor(*accent_rgb)
    _add_textbox(slide, Inches(3.82), Inches(2.58), Inches(1.56), Inches(0.7), center_label, size=18, bold=True, color=(255, 255, 255))

    positions = [
        (Inches(0.6), Inches(1.0)),
        (Inches(6.1), Inches(0.9)),
        (Inches(0.7), Inches(4.2)),
        (Inches(6.0), Inches(4.1)),
    ]
    for idx, node in enumerate(outer_nodes[:4]):
        left, top = positions[idx]
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, Inches(2.2), Inches(1.0))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*surface_rgb)
        shape.line.color.rgb = RGBColor(*accent_rgb)
        _add_textbox(slide, left + Inches(0.15), top + Inches(0.2), Inches(1.9), Inches(0.55), node, size=16, bold=True, color=text_rgb)
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            center.left + center.width // 2,
            center.top + center.height // 2,
            shape.left + shape.width // 2,
            shape.top + shape.height // 2,
        )
        connector.line.color.rgb = RGBColor(*accent_rgb)


def _render_table_slide(slide, slide_data, *, text_rgb):
    from pptx.util import Inches  # type: ignore

    table_data = slide_data.get("table") or {}
    headers = table_data.get("headers") or ["Column 1", "Column 2"]
    rows = table_data.get("rows") or []
    table_shape = slide.shapes.add_table(
        rows=len(rows) + 1,
        cols=len(headers),
        left=Inches(0.8),
        top=Inches(1.7),
        width=Inches(8.5),
        height=Inches(3.6),
    )
    table = table_shape.table
    for col_idx, header in enumerate(headers):
        table.cell(0, col_idx).text = str(header)
    for row_idx, row in enumerate(rows, start=1):
        for col_idx, value in enumerate(row[:len(headers)]):
            table.cell(row_idx, col_idx).text = str(value)
    if slide_data.get("summary"):
        _add_textbox(slide, Inches(0.85), Inches(5.55), Inches(8.0), Inches(0.5), slide_data["summary"], size=16, color=text_rgb)

def task_ppt_export(task_description: str, dependency_payload: str, **kwargs) -> str:
    try:
        from pptx import Presentation  # type: ignore
        from pptx.dml.color import RGBColor  # type: ignore
        from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE  # type: ignore
        from pptx.util import Inches  # type: ignore
    except ImportError:
        logger.error("python-pptx not installed. Run: pip install python-pptx")
        return "ERROR: python-pptx not installed"

    raw = invoke_task_llm(task_description + _PPT_PROMPT, dependency_payload, task_type="generation")

    try:
        import json5  # type: ignore
        data = json5.loads(raw)
    except Exception:
        import re
        match = re.search(r"\{.*\}", raw, re.DOTALL)
        data = json.loads(match.group(0)) if match else {"title": "Report", "slides": []}

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    theme = data.get("theme") or {}
    accent_rgb = _hex_to_rgb(theme.get("accent", "#1D4ED8"), (29, 78, 216))
    background_rgb = _hex_to_rgb(theme.get("background", "#F8FAFC"), (248, 250, 252))
    surface_rgb = _hex_to_rgb(theme.get("surface", "#E2E8F0"), (226, 232, 240))
    text_rgb = _hex_to_rgb(theme.get("text", "#0F172A"), (15, 23, 42))

    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = title_slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(*accent_rgb)
    title_box = _add_textbox(title_slide, Inches(0.9), Inches(1.0), Inches(8.9), Inches(1.2), data.get("title", "Report"), size=30, bold=True, color=(255, 255, 255))
    _add_textbox(title_slide, Inches(0.95), Inches(2.25), Inches(8.0), Inches(0.8), data.get("subtitle", "Strategic summary and connected thinking."), size=18, color=(235, 241, 255))
    if data.get("slides"):
        preview_points = [slide.get("title", "") for slide in data.get("slides", [])[:3]]
        _add_bullet_block(title_slide, Inches(0.95), Inches(3.35), Inches(5.2), Inches(2.2), preview_points, color=(255, 255, 255))
    hero = title_slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(9.1), Inches(1.1), Inches(3.1), Inches(4.9))
    hero.fill.solid()
    hero.fill.fore_color.rgb = RGBColor(*surface_rgb)
    hero.line.color.rgb = RGBColor(*surface_rgb)
    _add_textbox(title_slide, Inches(9.45), Inches(1.65), Inches(2.4), Inches(0.5), "Information Brain", size=18, bold=True, color=text_rgb)
    _add_textbox(title_slide, Inches(9.45), Inches(2.25), Inches(2.2), Inches(2.5), "This deck is structured to show the core answer, the connected themes behind it, and the brainstorm paths that come next.", size=15, color=text_rgb)

    for slide_data in data.get("slides", []):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(*background_rgb)
        _add_textbox(slide, Inches(0.7), Inches(0.45), Inches(10.8), Inches(0.55), slide_data.get("title", ""), size=24, bold=True, color=text_rgb)
        if slide_data.get("summary"):
            _add_textbox(slide, Inches(0.75), Inches(1.0), Inches(11.0), Inches(0.55), slide_data["summary"], size=16, color=text_rgb)

        layout = (slide_data.get("layout") or "summary").lower()
        if layout == "mind_map":
            _render_mind_map(slide, slide_data, accent_rgb=accent_rgb, surface_rgb=surface_rgb, text_rgb=text_rgb)
        elif layout == "brainstorm":
            ideas = slide_data.get("ideas") or slide_data.get("bullets", [])
            positions = [
                (Inches(0.8), Inches(1.8)),
                (Inches(3.2), Inches(1.5)),
                (Inches(5.7), Inches(2.0)),
                (Inches(8.2), Inches(1.7)),
            ]
            for idx, idea in enumerate(ideas[:4]):
                left, top = positions[idx]
                _add_card(slide, left, top, Inches(2.1), Inches(1.6), f"Idea {idx + 1}", idea, fill_rgb=surface_rgb, text_rgb=text_rgb)
        elif layout == "table":
            _render_table_slide(slide, slide_data, text_rgb=text_rgb)
        else:
            _add_bullet_block(slide, Inches(0.85), Inches(1.75), Inches(5.7), Inches(3.9), slide_data.get("bullets", []), color=text_rgb)
            if slide_data.get("ideas"):
                _add_card(slide, Inches(7.0), Inches(1.9), Inches(4.8), Inches(3.2), "Information Brain", "\n".join(slide_data.get("ideas", [])[:4]), fill_rgb=surface_rgb, text_rgb=text_rgb)

        if slide_data.get("notes"):
            slide.notes_slide.notes_text_frame.text = slide_data["notes"]

    path = os.path.join(EXPORT_ROOT, f"{uuid.uuid4().hex}.pptx")
    prs.save(path)
    logger.info(f"PPT saved: {path}")
    return path

"""PPTX parser — extracts slide text and speaker notes using python-pptx."""
from __future__ import annotations
import logging
logger = logging.getLogger(__name__)

def parse(file_path: str) -> str:
    try:
        from pptx import Presentation  # type: ignore
        prs = Presentation(file_path)
        slides = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
            notes = ""
            if slide.has_notes_slide:
                notes_tf = slide.notes_slide.notes_text_frame
                notes = notes_tf.text.strip()
            slide_text = "\n".join(texts)
            slides.append(f"--- Slide {i+1} ---\n{slide_text}" + (f"\n[Notes: {notes}]" if notes else ""))
        return "\n\n".join(slides)
    except Exception as e:
        return f"[PPTX PARSE ERROR: {e}]"
